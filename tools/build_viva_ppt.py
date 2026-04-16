from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = PROJECT_ROOT / "data" / "method_comparison_results.csv"
OUTPUT_PATH = PROJECT_ROOT / "Movie_Recommendation_LSH_Cosine_Viva.pptx"

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(30, 64, 175)
CYAN = RGBColor(14, 165, 233)
LIGHT_BG = RGBColor(244, 248, 252)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(31, 41, 55)
MUTED = RGBColor(75, 85, 99)


def add_fade_transition(slide) -> None:
    """Insert a subtle fade transition by editing slide XML."""
    sld = slide._element
    transition_tag = qn("p:transition")

    for existing in sld.findall(transition_tag):
        sld.remove(existing)

    transition = parse_xml(
        f"<p:transition {nsdecls('p')} advClick=\"1\"><p:fade/></p:transition>"
    )

    timing = sld.find(qn("p:timing"))
    if timing is not None:
        index = list(sld).index(timing)
        sld.insert(index, transition)
    else:
        sld.append(transition)


def set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_paragraph(
    paragraph, text: str, size: int, bold: bool, color: RGBColor, align=PP_ALIGN.LEFT
) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = align


def add_slide_number(slide, number: int, dark: bool = False) -> None:
    box = slide.shapes.add_textbox(
        Inches(12.65), Inches(7.05), Inches(0.5), Inches(0.3)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    style_paragraph(
        p,
        str(number),
        12,
        False,
        WHITE if dark else MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_top_banner(slide, title: str, subtitle: str, slide_number: int) -> None:
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(1.0),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0.95),
        Inches(13.333),
        Inches(0.05),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.18), Inches(8.6), Inches(0.4)
    )
    p_title = title_box.text_frame.paragraphs[0]
    p_title.clear()
    r_title = p_title.add_run()
    r_title.text = title
    r_title.font.name = FONT_TITLE
    r_title.font.size = Pt(24)
    r_title.font.bold = True
    r_title.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.57), Inches(8.8), Inches(0.28)
    )
    p_sub = subtitle_box.text_frame.paragraphs[0]
    style_paragraph(p_sub, subtitle, 12, False, RGBColor(191, 219, 254))

    add_slide_number(slide, slide_number)


def add_bullets(
    shape, lines: list[str], font_size: int = 20, color: RGBColor = TEXT_DARK
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        style_paragraph(p, line, font_size, False, color)
        p.level = 0
        p.space_after = Pt(8)


def build_presentation() -> Path:
    df = pd.read_csv(DATA_PATH)

    avg_cosine = float(df["cosine_time_ms"].mean())
    avg_lsh = float(df["lsh_time_ms"].mean())
    speedup_pct = ((avg_cosine - avg_lsh) / avg_cosine) * 100
    lsh_faster_count = int((df["lsh_time_ms"] < df["cosine_time_ms"]).sum())
    cosine_faster_count = int(len(df) - lsh_faster_count)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1: Title
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, NAVY)

    decorative_1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(10.2),
        Inches(-0.6),
        Inches(4.2),
        Inches(4.2),
    )
    decorative_1.fill.solid()
    decorative_1.fill.fore_color.rgb = BLUE
    decorative_1.fill.transparency = 0.3
    decorative_1.line.fill.background()

    decorative_2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-1.5),
        Inches(5.0),
        Inches(4.8),
        Inches(4.8),
    )
    decorative_2.fill.solid()
    decorative_2.fill.fore_color.rgb = CYAN
    decorative_2.fill.transparency = 0.35
    decorative_2.line.fill.background()

    label_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.6), Inches(3.8), Inches(0.4)
    )
    style_paragraph(
        label_box.text_frame.paragraphs[0],
        "Final Year Project Viva",
        15,
        True,
        RGBColor(147, 197, 253),
    )

    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(10.3), Inches(1.9)
    )
    tf_title = title_box.text_frame
    tf_title.clear()
    p1 = tf_title.paragraphs[0]
    p1.clear()
    r1 = p1.add_run()
    r1.text = "Movie Recommendation System using\nLSH and Cosine Similarity"
    r1.font.name = FONT_TITLE
    r1.font.size = Pt(44)
    r1.font.bold = True
    r1.font.color.rgb = WHITE

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.7), Inches(8.5), Inches(0.6)
    )
    style_paragraph(
        subtitle_box.text_frame.paragraphs[0],
        "Scalable Personalized Recommendation Approach",
        24,
        False,
        RGBColor(191, 219, 254),
    )

    info_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(4.65),
        Inches(8.4),
        Inches(2.0),
    )
    info_card.fill.solid()
    info_card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    info_card.fill.transparency = 0.18
    info_card.line.color.rgb = RGBColor(125, 211, 252)
    info_card.line.width = Pt(1.25)

    placeholders = slide.shapes.add_textbox(
        Inches(1.1), Inches(5.0), Inches(7.8), Inches(1.4)
    )
    add_bullets(
        placeholders,
        [
            "Name: _________________________",
            "Roll Number: __________________",
            "College / Department: __________________",
        ],
        font_size=18,
        color=WHITE,
    )

    add_slide_number(slide, 1, dark=True)
    add_fade_transition(slide)
    add_notes(
        slide,
        "Open with project motivation: users face content overload on streaming platforms.\n"
        "State that this project compares two recommendation paradigms: Cosine Similarity and LSH.\n"
        "Mention that LSH is selected for scalable serving while cosine is used as a quality baseline.\n"
        "Introduce your identity details from placeholders before moving to background.",
    )

    # Slide 2: Introduction
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Introduction",
        "Why recommendation systems matter in modern content platforms",
        2,
    )

    intro_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.35), Inches(7.25), Inches(4.8)
    )
    add_bullets(
        intro_box,
        [
            "A movie recommendation system suggests relevant movies based on user preferences and content similarity.",
            "Personalization reduces decision fatigue and increases user engagement.",
            "It helps platforms improve retention, watch time, and overall user satisfaction.",
        ],
        font_size=21,
    )

    apps_title = slide.shapes.add_textbox(
        Inches(8.45), Inches(1.45), Inches(4.2), Inches(0.35)
    )
    style_paragraph(
        apps_title.text_frame.paragraphs[0], "Real-World Applications", 16, True, BLUE
    )

    app_cards = [
        ("Netflix", "Personalized home feed and next-watch suggestions"),
        ("Prime Video", "Context-aware and genre-based recommendations"),
        ("YouTube", "Session-level recommendation and ranking"),
    ]
    y = 1.9
    for app_name, app_desc in app_cards:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.45),
            Inches(y),
            Inches(4.0),
            Inches(1.35),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(203, 213, 225)

        app_name_box = slide.shapes.add_textbox(
            Inches(8.7), Inches(y + 0.2), Inches(1.4), Inches(0.3)
        )
        style_paragraph(app_name_box.text_frame.paragraphs[0], app_name, 15, True, NAVY)

        app_desc_box = slide.shapes.add_textbox(
            Inches(8.7), Inches(y + 0.55), Inches(3.55), Inches(0.65)
        )
        style_paragraph(
            app_desc_box.text_frame.paragraphs[0], app_desc, 12, False, MUTED
        )

        y += 1.5

    add_fade_transition(slide)
    add_notes(
        slide,
        "Define recommendation system in one sentence before discussing significance.\n"
        "Emphasize that personalization directly affects business KPIs like retention.\n"
        "Use Netflix, Prime Video, and YouTube as familiar examples for faculty.\n"
        "Conclude this slide by linking motivation to scalability challenges in the next slide.",
    )

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Problem Statement",
        "Key bottlenecks in traditional recommendation pipelines",
        3,
    )

    challenge_title = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.35), Inches(4.5), Inches(0.4)
    )
    style_paragraph(
        challenge_title.text_frame.paragraphs[0], "Challenges", 20, True, NAVY
    )

    challenges = [
        (
            "Large Dataset",
            "High-dimensional metadata and thousands of movies increase search space.",
        ),
        (
            "Slow Similarity Search",
            "Pairwise similarity checks become expensive for real-time serving.",
        ),
        (
            "Scalability Issue",
            "Latency grows as data volume and concurrent users increase.",
        ),
    ]

    x = 0.8
    for title, desc in challenges:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.9),
            Inches(3.95),
            Inches(2.35),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = RGBColor(203, 213, 225)

        title_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(2.1), Inches(3.5), Inches(0.35)
        )
        style_paragraph(title_box.text_frame.paragraphs[0], title, 16, True, BLUE)

        desc_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(2.55), Inches(3.5), Inches(1.35)
        )
        style_paragraph(desc_box.text_frame.paragraphs[0], desc, 13, False, TEXT_DARK)

        x += 4.2

    objective_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(4.75),
        Inches(11.9),
        Inches(1.75),
    )
    objective_box.fill.solid()
    objective_box.fill.fore_color.rgb = RGBColor(219, 234, 254)
    objective_box.line.color.rgb = RGBColor(147, 197, 253)

    objective_title = slide.shapes.add_textbox(
        Inches(1.1), Inches(5.05), Inches(2.4), Inches(0.35)
    )
    style_paragraph(
        objective_title.text_frame.paragraphs[0], "Project Objective", 16, True, NAVY
    )

    objective_desc = slide.shapes.add_textbox(
        Inches(1.1), Inches(5.45), Inches(11.2), Inches(0.75)
    )
    style_paragraph(
        objective_desc.text_frame.paragraphs[0],
        "Design a recommendation pipeline that delivers fast response time with acceptable recommendation quality by comparing LSH against cosine similarity.",
        15,
        False,
        TEXT_DARK,
    )

    add_fade_transition(slide)
    add_notes(
        slide,
        "State that brute-force similarity is the core bottleneck.\n"
        "Discuss why latency and scalability are critical for production recommendation systems.\n"
        "Mention that this project targets a practical trade-off: speed and quality.\n"
        "Transition to LSH as the method used to address these bottlenecks.",
    )

    # Slide 4: Methodology - LSH
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Methodology - Locality Sensitive Hashing (LSH)",
        "Approximate nearest-neighbor search for scalable recommendation",
        4,
    )

    method_points = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.35), Inches(12.0), Inches(1.1)
    )
    add_bullets(
        method_points,
        [
            "LSH groups similar movie representations into the same hash buckets.",
            "Instead of searching all movies, we search only relevant buckets for fast retrieval.",
        ],
        font_size=18,
    )

    stages = [
        "Dataset",
        "Preprocessing",
        "Hashing",
        "Similar Bucket\nSearch",
        "Recommendation",
    ]

    x = 0.65
    y = 3.1
    box_w = 2.25
    box_h = 1.35

    for i, stage in enumerate(stages):
        stage_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(box_w),
            Inches(box_h),
        )
        stage_box.fill.solid()
        stage_box.fill.fore_color.rgb = WHITE
        stage_box.line.color.rgb = CYAN
        stage_box.line.width = Pt(1.3)

        label = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.35), Inches(box_w - 0.3), Inches(0.7)
        )
        p = label.text_frame.paragraphs[0]
        style_paragraph(p, stage, 14, True, NAVY, align=PP_ALIGN.CENTER)
        label.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + box_w + 0.05),
                Inches(y + 0.5),
                Inches(0.35),
                Inches(0.35),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.fill.background()

        x += 2.62

    note_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(5.0),
        Inches(12.0),
        Inches(1.25),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RGBColor(236, 253, 245)
    note_box.line.color.rgb = RGBColor(134, 239, 172)

    note_text = slide.shapes.add_textbox(
        Inches(1.05), Inches(5.35), Inches(11.4), Inches(0.6)
    )
    style_paragraph(
        note_text.text_frame.paragraphs[0],
        "Implementation uses MinHash signatures and LSH indexing to reduce search complexity while preserving semantic similarity.",
        14,
        False,
        TEXT_DARK,
    )

    add_fade_transition(slide)
    add_notes(
        slide,
        "Explain LSH in simple words: similar items collide in the same bucket with high probability.\n"
        "Walk through each step in the flowchart from data to recommendation output.\n"
        "Mention that this avoids full pairwise comparisons and improves response time.\n"
        "Briefly connect this to your implementation using MinHash + LSH index.",
    )

    # Slide 5: Cosine Similarity
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Cosine Similarity Approach",
        "Vector-angle based similarity for content recommendation",
        5,
    )

    cosine_desc = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(7.3), Inches(1.8)
    )
    add_bullets(
        cosine_desc,
        [
            "Cosine similarity measures how close two movie vectors are based on the angle between them.",
            "Higher cosine value means stronger similarity in movie content features.",
        ],
        font_size=18,
    )

    formula_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(3.35),
        Inches(6.95),
        Inches(1.55),
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = WHITE
    formula_box.line.color.rgb = RGBColor(147, 197, 253)

    formula_title = slide.shapes.add_textbox(
        Inches(1.2), Inches(3.55), Inches(1.8), Inches(0.3)
    )
    style_paragraph(formula_title.text_frame.paragraphs[0], "Formula", 14, True, BLUE)

    formula = slide.shapes.add_textbox(
        Inches(1.2), Inches(3.95), Inches(6.4), Inches(0.7)
    )
    p_formula = formula.text_frame.paragraphs[0]
    style_paragraph(
        p_formula,
        "cos(theta) = (A . B) / (||A|| x ||B||)",
        26,
        True,
        NAVY,
        align=PP_ALIGN.CENTER,
    )

    strengths = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.35),
        Inches(1.8),
        Inches(4.15),
        Inches(1.65),
    )
    strengths.fill.solid()
    strengths.fill.fore_color.rgb = WHITE
    strengths.line.color.rgb = RGBColor(203, 213, 225)

    strengths_text = slide.shapes.add_textbox(
        Inches(8.6), Inches(2.05), Inches(3.7), Inches(1.1)
    )
    add_bullets(
        strengths_text,
        [
            "Accurate for small to medium datasets",
            "Strong baseline for similarity quality",
        ],
        font_size=13,
        color=TEXT_DARK,
    )

    limitations = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.35),
        Inches(3.75),
        Inches(4.15),
        Inches(1.65),
    )
    limitations.fill.solid()
    limitations.fill.fore_color.rgb = RGBColor(254, 242, 242)
    limitations.line.color.rgb = RGBColor(252, 165, 165)

    limitations_text = slide.shapes.add_textbox(
        Inches(8.6), Inches(4.0), Inches(3.75), Inches(1.15)
    )
    add_bullets(
        limitations_text,
        [
            "Computationally expensive for large data",
            "Response time increases with dataset size",
        ],
        font_size=13,
        color=TEXT_DARK,
    )

    add_fade_transition(slide)
    add_notes(
        slide,
        "Describe cosine similarity as a geometric measure between TF-IDF vectors.\n"
        "Read the formula once and explain each term in simple language.\n"
        "Mention that cosine gives strong quality but scales poorly for exhaustive search.\n"
        "Set up the transition to quantitative comparison in the next slide.",
    )

    # Slide 6: Results and Comparison
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Results and Comparison",
        "Empirical runtime and qualitative quality comparison",
        6,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["Cosine Similarity", "LSH"]
    chart_data.add_series(
        "Average Query Time (ms)", [round(avg_cosine, 2), round(avg_lsh, 2)]
    )

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(6.35),
        Inches(3.85),
        chart_data,
    ).chart

    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.name = FONT_BODY
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.name = FONT_BODY
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = "Average Time per Query"
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = FONT_TITLE
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

    series = chart.series[0]
    series.has_data_labels = True
    series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    series.data_labels.number_format = "0.00"
    series.data_labels.font.size = Pt(10)

    point_cos = series.points[0]
    point_cos.format.fill.solid()
    point_cos.format.fill.fore_color.rgb = BLUE

    point_lsh = series.points[1]
    point_lsh.format.fill.solid()
    point_lsh.format.fill.fore_color.rgb = CYAN

    table_shape = slide.shapes.add_table(
        5,
        3,
        Inches(7.35),
        Inches(1.5),
        Inches(5.35),
        Inches(3.85),
    )
    table = table_shape.table
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(1.55)
    table.columns[2].width = Inches(1.55)

    table_data = [
        ("Metric", "Cosine", "LSH"),
        ("Recommendation Accuracy", "High", "Good"),
        ("Avg Time (ms)", f"{avg_cosine:.2f}", f"{avg_lsh:.2f}"),
        ("Scalability", "Moderate", "High"),
        ("Faster Query Wins", f"{cosine_faster_count}/20", f"{lsh_faster_count}/20"),
    ]

    for r in range(5):
        for c in range(3):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if not p.runs:
                p.add_run()
            run = p.runs[0]
            run.font.name = FONT_BODY
            run.font.size = Pt(11)
            run.font.bold = r == 0
            run.font.color.rgb = WHITE if r == 0 else TEXT_DARK

    highlight = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(5.55),
        Inches(11.9),
        Inches(1.1),
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RGBColor(220, 252, 231)
    highlight.line.color.rgb = RGBColor(134, 239, 172)

    highlight_text = slide.shapes.add_textbox(
        Inches(1.05), Inches(5.86), Inches(11.4), Inches(0.55)
    )
    style_paragraph(
        highlight_text.text_frame.paragraphs[0],
        f"Key Insight: LSH is {speedup_pct:.1f}% faster on average, while cosine remains slightly more precise for fine-grained ranking.",
        15,
        True,
        RGBColor(22, 101, 52),
        align=PP_ALIGN.CENTER,
    )

    add_fade_transition(slide)
    add_notes(
        slide,
        f"Present measured timing from comparison CSV: cosine {avg_cosine:.2f} ms vs LSH {avg_lsh:.2f} ms.\n"
        f"Mention LSH wins {lsh_faster_count} out of 20 query cases, indicating better runtime scalability.\n"
        "State accuracy as a qualitative comparison: cosine is more exact, LSH is near-accurate with speed advantage.\n"
        "Conclude with the trade-off: choose method based on scale and latency requirement.",
    )

    # Slide 7: Conclusion and Future Scope
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, LIGHT_BG)
    add_top_banner(
        slide,
        "Conclusion and Future Scope",
        "Project takeaways and potential enhancements",
        7,
    )

    conclusion_title = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.4), Inches(5.0), Inches(0.4)
    )
    style_paragraph(
        conclusion_title.text_frame.paragraphs[0], "Conclusion", 20, True, NAVY
    )

    conclusion_points = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.9), Inches(6.3), Inches(3.8)
    )
    add_bullets(
        conclusion_points,
        [
            "LSH improves efficiency for large datasets by reducing search space.",
            "Cosine similarity remains a strong baseline for recommendation quality.",
            "The project demonstrates a practical balance between speed and relevance.",
        ],
        font_size=18,
    )

    future_title = slide.shapes.add_textbox(
        Inches(7.2), Inches(1.4), Inches(5.2), Inches(0.4)
    )
    style_paragraph(
        future_title.text_frame.paragraphs[0], "Future Scope", 20, True, NAVY
    )

    futures = [
        "Hybrid recommender system (LSH + cosine reranking)",
        "Deep learning based embeddings for richer semantics",
        "Real-time recommendation with streaming user behavior",
    ]

    y = 1.95
    for item in futures:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(7.2),
            Inches(y),
            Inches(5.1),
            Inches(1.1),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(203, 213, 225)

        item_box = slide.shapes.add_textbox(
            Inches(7.45), Inches(y + 0.28), Inches(4.6), Inches(0.55)
        )
        style_paragraph(item_box.text_frame.paragraphs[0], item, 13, False, TEXT_DARK)
        y += 1.3

    thanks = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(6.1),
        Inches(11.5),
        Inches(0.85),
    )
    thanks.fill.solid()
    thanks.fill.fore_color.rgb = NAVY
    thanks.line.fill.background()

    thanks_text = slide.shapes.add_textbox(
        Inches(1.0), Inches(6.3), Inches(11.1), Inches(0.4)
    )
    style_paragraph(
        thanks_text.text_frame.paragraphs[0],
        "Thank You - Questions and Discussion",
        18,
        True,
        WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_fade_transition(slide)
    add_notes(
        slide,
        "Summarize that LSH addresses scalability while preserving useful recommendation quality.\n"
        "Reinforce cosine as a baseline for benchmarking and validation.\n"
        "Present future work in three tracks: hybrid modeling, deep embeddings, and real-time serving.\n"
        "Close confidently and invite questions from the panel.",
    )

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output_file = build_presentation()
    print(f"Presentation created: {output_file}")
