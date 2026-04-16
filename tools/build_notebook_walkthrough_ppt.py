from __future__ import annotations

import argparse
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"
FONT_CODE = "Consolas"

NAVY = RGBColor(16, 33, 63)
BLUE = RGBColor(41, 98, 255)
LIGHT_BG = RGBColor(246, 248, 252)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(32, 40, 52)
MUTED = RGBColor(90, 102, 120)
BORDER = RGBColor(219, 224, 232)


ACTION_PATTERNS: list[tuple[str, str]] = [
    (
        r"\bread_csv\b|\bread_excel\b|\bread_parquet\b",
        "Dataset files load kiye ja rahe hain.",
    ),
    (r"\bmerge\b|\bjoin\b", "Multiple datasets combine/merge ho rahe hain."),
    (
        r"\bshape\b|\bhead\b|\binfo\b|\bdescribe\b|\bisnull\b",
        "Data ka inspection/EDA snapshot nikala ja raha hai.",
    ),
    (
        r"\bdropna\b|\bdrop_duplicates\b",
        "Missing ya duplicate rows cleanup ki ja rahi hai.",
    ),
    (
        r"literal_eval|ast\.|json",
        "Text-based JSON/list values parse kiye ja rahe hain.",
    ),
    (
        r"split\(|join\(|lower\(|replace\(|stem|token",
        "Text preprocessing aur feature normalization ho raha hai.",
    ),
    (r"apply\(|lambda", "Row-wise transformation/feature engineering run ho rahi hai."),
    (
        r"to_csv\(|to_pickle\(|dump\(",
        "Processed artifact/output file save kiya ja raha hai.",
    ),
]

WORKFLOW_PATTERNS: list[tuple[str, str]] = [
    ("Data files load", r"\bread_csv\b|\bread_excel\b|\bread_parquet\b"),
    ("Initial data inspection", r"\bshape\b|\bhead\b|\binfo\b|\bdescribe\b|\bisnull\b"),
    ("Dataset merge/join", r"\bmerge\b|\bjoin\b"),
    ("Column filtering/selection", r"\[\[|\bdrop\(|\brename\("),
    ("Missing and duplicate cleanup", r"\bdropna\b|\bdrop_duplicates\b"),
    ("JSON/text parsing", r"literal_eval|ast\.|json"),
    (
        "Tag/text feature engineering",
        r"apply\(|lambda|split\(|join\(|lower\(|replace\(|stem",
    ),
    ("Output export", r"to_csv\(|to_pickle\(|dump\("),
]


def _to_text(value: Any) -> str:
    if isinstance(value, list):
        return "".join(str(item) for item in value)
    return str(value)


def _collapse_spaces(line: str) -> str:
    return re.sub(r"[ \t]+", " ", line).strip()


def clean_multiline_text(text: str) -> str:
    lines = [_collapse_spaces(line) for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines)


def clean_markdown(text: str) -> str:
    text = text.replace("```", "")
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", text)
    text = re.sub(r"^\s*[-*]\s+", "- ", text, flags=re.MULTILINE)
    return clean_multiline_text(text)


def trim_for_slide(text: str, max_lines: int = 22, max_chars: int = 1400) -> str:
    lines = text.splitlines()
    if len(lines) > max_lines:
        lines = lines[:max_lines] + ["... (trimmed)"]
    result = "\n".join(lines)
    if len(result) > max_chars:
        result = result[: max_chars - 12] + "... (trimmed)"
    return result


def extract_output_snapshot(outputs: list[dict[str, Any]]) -> str:
    if not outputs:
        return "No visible output in this cell."

    chunks: list[str] = []
    image_count = 0

    for output in outputs:
        output_type = output.get("output_type", "")

        if output_type == "stream":
            stream_text = _to_text(output.get("text", "")).strip()
            if stream_text:
                chunks.append(stream_text)

        elif output_type in {"execute_result", "display_data"}:
            data = output.get("data", {})

            if "text/plain" in data:
                plain_text = _to_text(data.get("text/plain", "")).strip()
                if plain_text:
                    chunks.append(plain_text)

            if "text/html" in data and "text/plain" not in data:
                chunks.append("HTML output generated (table/formatting).")

            if "image/png" in data or "image/jpeg" in data:
                image_count += 1

        elif output_type == "error":
            ename = output.get("ename", "Error")
            evalue = output.get("evalue", "")
            chunks.append(f"{ename}: {evalue}".strip())

    if image_count:
        chunks.append(f"Image/chart output generated: {image_count}")

    if not chunks:
        return "Output present but not text-readable in this export."

    return trim_for_slide(
        clean_multiline_text("\n\n".join(chunks)), max_lines=18, max_chars=1100
    )


def infer_actions_from_code(code: str) -> list[str]:
    if not code.strip():
        return ["Empty code cell."]

    code_lower = code.lower()
    actions: list[str] = []

    for pattern, message in ACTION_PATTERNS:
        if re.search(pattern, code_lower):
            actions.append(message)

    if not actions:
        actions.append("Is cell me custom logic/calculation execute ho rahi hai.")

    # keep concise for slide readability
    return actions[:4]


def detect_workflow_steps(all_code: str) -> list[str]:
    steps: list[str] = []
    code_lower = all_code.lower()

    for label, pattern in WORKFLOW_PATTERNS:
        if re.search(pattern, code_lower):
            steps.append(label)

    if not steps:
        steps.append("Notebook me custom transformation workflow hai.")

    return steps


def detect_files_from_code(all_code: str) -> list[str]:
    files: list[str] = []
    for match in re.finditer(
        r"[\"']([^\"']+\.(?:csv|pkl|pickle|json))[\"']", all_code, flags=re.IGNORECASE
    ):
        file_path = match.group(1)
        if file_path not in files:
            files.append(file_path)
    return files[:10]


def style_paragraph(
    paragraph,
    text: str,
    size: int,
    *,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT_BODY,
) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = align


def add_header(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(0.95),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.15), Inches(9.5), Inches(0.4)
    )
    p_title = title_box.text_frame.paragraphs[0]
    style_paragraph(p_title, title, 24, bold=True, color=WHITE, font_name=FONT_TITLE)

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(0.56), Inches(12.2), Inches(0.25)
    )
    p_subtitle = subtitle_box.text_frame.paragraphs[0]
    style_paragraph(
        p_subtitle, subtitle, 12, color=RGBColor(198, 214, 246), font_name=FONT_BODY
    )


def add_content_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    heading: str,
    body: str,
    monospace: bool = False,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER

    heading_box = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.3)
    )
    p_heading = heading_box.text_frame.paragraphs[0]
    style_paragraph(p_heading, heading, 14, bold=True, color=BLUE)

    body_box = slide.shapes.add_textbox(
        Inches(x + 0.2), Inches(y + 0.5), Inches(w - 0.4), Inches(h - 0.62)
    )
    tf = body_box.text_frame
    tf.clear()
    tf.word_wrap = True

    lines = body.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        style_paragraph(
            paragraph,
            line,
            11 if monospace else 13,
            color=TEXT_DARK,
            font_name=FONT_CODE if monospace else FONT_BODY,
        )


def add_title_slide(prs: Presentation, notebook_path: Path, cell_count: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    add_header(
        slide,
        "Notebook Walkthrough PPT",
        f"Source: {notebook_path.name}",
    )

    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.6), Inches(12.0), Inches(1.6)
    )
    tf_title = title_box.text_frame
    p1 = tf_title.paragraphs[0]
    style_paragraph(
        p1,
        notebook_path.stem.replace("_", " ").title(),
        40,
        bold=True,
        color=NAVY,
        font_name=FONT_TITLE,
    )

    info_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.5), Inches(12.0), Inches(1.0)
    )
    p2 = info_box.text_frame.paragraphs[0]
    generated_on = datetime.now().strftime("%d %b %Y, %I:%M %p")
    style_paragraph(
        p2,
        f"Cell-wise explanation and outputs | Total cells: {cell_count} | Generated: {generated_on}",
        17,
        color=MUTED,
    )


def add_overview_slide(
    prs: Presentation,
    notebook_name: str,
    total_cells: int,
    code_cells: int,
    markdown_cells: int,
    output_cells: int,
    workflow_steps: list[str],
    detected_files: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    add_header(slide, "Notebook Summary", "Kya ho raha hai: high-level view")

    left_body_lines = [
        f"- Notebook: {notebook_name}",
        f"- Total cells: {total_cells}",
        f"- Code cells: {code_cells}",
        f"- Markdown cells: {markdown_cells}",
        f"- Code cells with output: {output_cells}",
        "",
        "Detected workflow:",
    ]
    left_body_lines.extend([f"- {step}" for step in workflow_steps])

    right_body_lines = ["Detected input/output files from code:"]
    if detected_files:
        right_body_lines.extend([f"- {file_path}" for file_path in detected_files])
    else:
        right_body_lines.append("- No explicit file paths detected.")

    add_content_card(
        slide,
        x=0.4,
        y=1.15,
        w=7.4,
        h=5.9,
        heading="Overview",
        body="\n".join(left_body_lines),
    )

    add_content_card(
        slide,
        x=8.0,
        y=1.15,
        w=4.9,
        h=5.9,
        heading="Files",
        body="\n".join(right_body_lines),
    )


def add_cell_slide(
    prs: Presentation,
    cell_number: int,
    cell_type: str,
    source_text: str,
    action_notes: list[str],
    output_snapshot: str,
    execution_count: Any,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG

    exec_label = "N/A" if execution_count is None else str(execution_count)
    add_header(
        slide,
        f"Cell {cell_number}: {cell_type.title()}",
        f"Execution count: {exec_label}",
    )

    left_heading = "Markdown Content" if cell_type == "markdown" else "Code"
    left_body = trim_for_slide(source_text, max_lines=24, max_chars=1700)

    action_body = "\n".join([f"- {note}" for note in action_notes])

    add_content_card(
        slide,
        x=0.4,
        y=1.15,
        w=7.8,
        h=5.95,
        heading=left_heading,
        body=left_body,
        monospace=cell_type == "code",
    )

    add_content_card(
        slide,
        x=8.35,
        y=1.15,
        w=4.55,
        h=2.75,
        heading="Kya ho raha hai",
        body=action_body,
    )

    add_content_card(
        slide,
        x=8.35,
        y=4.0,
        w=4.55,
        h=3.1,
        heading="Output Snapshot",
        body=trim_for_slide(output_snapshot, max_lines=16, max_chars=900),
        monospace=False,
    )


def build_walkthrough_ppt(notebook_path: Path, output_path: Path) -> tuple[Path, int]:
    with notebook_path.open("r", encoding="utf-8") as file:
        notebook = json.load(file)

    cells = notebook.get("cells", [])
    total_cells = len(cells)
    code_cells = [cell for cell in cells if cell.get("cell_type") == "code"]
    markdown_cells = [cell for cell in cells if cell.get("cell_type") == "markdown"]

    output_cells = 0
    all_code_text_parts: list[str] = []

    for cell in code_cells:
        source = _to_text(cell.get("source", ""))
        all_code_text_parts.append(source)
        outputs = cell.get("outputs", [])
        if outputs:
            output_cells += 1

    all_code_text = "\n".join(all_code_text_parts)
    workflow_steps = detect_workflow_steps(all_code_text)
    detected_files = detect_files_from_code(all_code_text)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs, notebook_path, total_cells)

    add_overview_slide(
        prs,
        notebook_name=notebook_path.name,
        total_cells=total_cells,
        code_cells=len(code_cells),
        markdown_cells=len(markdown_cells),
        output_cells=output_cells,
        workflow_steps=workflow_steps,
        detected_files=detected_files,
    )

    for index, cell in enumerate(cells, start=1):
        cell_type = cell.get("cell_type", "unknown")
        source_raw = _to_text(cell.get("source", ""))

        if cell_type == "markdown":
            source_text = clean_markdown(source_raw)
            action_notes = [
                "Section heading / narrative explanation di gayi hai.",
                "Yeh agle code cells ka context set karta hai.",
            ]
            output_snapshot = "Markdown cell me executable output nahi hota."
            execution_count = None
        else:
            source_text = (
                source_raw.strip() if source_raw.strip() else "(empty code cell)"
            )
            action_notes = infer_actions_from_code(source_raw)
            output_snapshot = extract_output_snapshot(cell.get("outputs", []))
            execution_count = cell.get("execution_count")

        add_cell_slide(
            prs,
            cell_number=index,
            cell_type=cell_type,
            source_text=source_text,
            action_notes=action_notes,
            output_snapshot=output_snapshot,
            execution_count=execution_count,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path, len(prs.slides)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a full cell-by-cell walkthrough PPT from a Jupyter notebook."
    )
    parser.add_argument(
        "--notebook",
        type=Path,
        default=Path("notebooks/01_preprocessing.ipynb"),
        help="Path to input notebook file.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("notebooks/01_preprocessing_walkthrough.pptx"),
        help="Path for generated PPTX file.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    notebook_path = args.notebook.resolve()
    output_path = args.output.resolve()

    if not notebook_path.exists():
        raise FileNotFoundError(f"Notebook not found: {notebook_path}")

    saved_path, slide_count = build_walkthrough_ppt(notebook_path, output_path)
    print(f"Notebook walkthrough PPT generated: {saved_path}")
    print(f"Total slides: {slide_count}")


if __name__ == "__main__":
    main()
